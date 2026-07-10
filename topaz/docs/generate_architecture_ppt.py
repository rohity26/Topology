#!/usr/bin/env python3
"""
Generate O-RAN Network Architecture PPT
========================================
Creates a PowerPoint presentation documenting the dual-network
architecture (oran-intel + telemetry-net) and tc packet loss behavior.

Usage:
    python3 docs/generate_architecture_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "O-RAN_Network_Architecture.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x16, 0x21, 0x3E)
BLUE = RGBColor(0x00, 0xD2, 0xFF)
GREEN = RGBColor(0x00, 0xFF, 0x88)
RED = RGBColor(0xFF, 0x44, 0x44)
ORANGE = RGBColor(0xFF, 0xAA, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE = RGBColor(0xBB, 0x86, 0xFC)
YELLOW = RGBColor(0xFF, 0xEB, 0x3B)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=12,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Consolas"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None,
             text="", font_size=10, font_color=WHITE):
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = "Consolas"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape


# =========================================================================
# SLIDE 1: Title
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 1, 1.5, 11, 1,
            "O-RAN srsRAN FALCON Extension",
            font_size=36, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 2.8, 11, 1,
            "Dual-Network Architecture: Data Plane & Metrics Separation",
            font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 4.2, 11, 0.8,
            "175.0.0.0/8 (oran-intel)  -  Data Plane (F1AP / F1U / N2 / ZMQ)",
            font_size=16, color=GREEN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 4.9, 11, 0.8,
            "171.0.0.0/8 (telemetry-net)  -  Metrics Collection (UDP → InfluxDB → Grafana)",
            font_size=16, color=ORANGE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 6.0, 11, 0.5,
            "Setup Commands:\n"
            "docker network create --subnet=175.0.0.0/8 oran-intel\n"
            "docker network create --subnet=171.0.0.0/8 telemetry-net",
            font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 2: Full Network Topology
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.3, 0.2, 12, 0.6,
            "Network Topology - Dual Subnet Architecture",
            font_size=24, color=BLUE, bold=True)

# oran-intel box
add_rect(slide, 0.3, 0.9, 6.2, 6.0, BG_CARD, GREEN)
add_textbox(slide, 0.5, 0.95, 5.8, 0.4,
            "175.0.0.0/8 (oran-intel) - DATA PLANE",
            font_size=13, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# 5GC
add_rect(slide, 2.5, 1.5, 1.8, 0.55, RGBColor(0x00, 0x44, 0x00), GREEN,
         "5GC\n175.53.1.2", font_size=9, font_color=GREEN)

# CUs
cu_data = [("cu0", "175.53.10.1\n2nd: .1.11", 0.7),
           ("cu1", "175.53.10.2\n2nd: .1.12", 2.55),
           ("cu2", "175.53.10.3\n2nd: .1.15", 4.4)]
for name, ip, left in cu_data:
    add_rect(slide, left, 2.4, 1.6, 0.7, RGBColor(0x00, 0x22, 0x55), BLUE,
             f"{name}\n{ip}", font_size=8, font_color=BLUE)

# DUs
du_data = [("du0\n.1.10", 0.4), ("du1\n.1.13", 1.35),
           ("du2\n.1.14", 2.3), ("du3\n.1.16", 3.25),
           ("du4\n.1.17", 4.2), ("du5\n.1.18", 5.15)]
for label, left in du_data:
    add_rect(slide, left, 3.5, 0.85, 0.6, RGBColor(0x22, 0x33, 0x55), BLUE,
             label, font_size=8, font_color=WHITE)

# UEs
ue_data = [("ue0\n.2.1", 0.4), ("ue1\n.3.1", 1.35),
           ("ue2\n.4.1", 2.3), ("ue3\n.5.1", 3.25),
           ("ue4\n.6.1", 4.2), ("ue5\n.7.1", 5.15)]
for label, left in ue_data:
    add_rect(slide, left, 4.5, 0.85, 0.6, RGBColor(0x33, 0x22, 0x55), PURPLE,
             label, font_size=8, font_color=WHITE)

add_textbox(slide, 0.5, 5.3, 5.5, 0.3,
            "Redis: 175.24.0.1",
            font_size=10, color=GRAY)

add_textbox(slide, 0.5, 5.7, 5.5, 1.0,
            "Protocols: N2(SCTP:38412), F1AP(SCTP:38472),\n"
            "F1U(UDP), ZMQ(TCP) between DU<->UE",
            font_size=10, color=GRAY)

# telemetry-net box
add_rect(slide, 6.8, 0.9, 6.2, 6.0, BG_CARD, ORANGE)
add_textbox(slide, 7.0, 0.95, 5.8, 0.4,
            "171.0.0.0/8 (telemetry-net) - METRICS PLANE",
            font_size=13, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Metrics infra
add_rect(slide, 7.2, 1.5, 1.6, 0.55, RGBColor(0x44, 0x22, 0x00), ORANGE,
         "metrics-server\n171.40.1.4:55555", font_size=8, font_color=ORANGE)
add_rect(slide, 9.0, 1.5, 1.4, 0.55, RGBColor(0x44, 0x22, 0x00), ORANGE,
         "influxdb\n171.40.1.5:8086", font_size=8, font_color=ORANGE)
add_rect(slide, 10.6, 1.5, 1.4, 0.55, RGBColor(0x44, 0x22, 0x00), ORANGE,
         "grafana\n171.40.1.6:3000", font_size=8, font_color=ORANGE)

# CUs on telemetry-net
cu_met = [("cu0\n171.53.10.1", 7.2), ("cu1\n171.53.10.2", 8.8),
          ("cu2\n171.53.10.3", 10.4)]
for label, left in cu_met:
    add_rect(slide, left, 2.5, 1.4, 0.55, RGBColor(0x00, 0x22, 0x33), ORANGE,
             label, font_size=8, font_color=WHITE)

# DUs on telemetry-net
du_met = [("du0\n171.53.1.10", 7.0), ("du1\n171.53.1.13", 8.0),
          ("du2\n171.53.1.14", 9.0), ("du3\n171.53.1.16", 10.0),
          ("du4\n171.53.1.17", 11.0), ("du5\n171.53.1.18", 12.0)]
for label, left in du_met:
    add_rect(slide, left, 3.5, 0.9, 0.6, RGBColor(0x33, 0x22, 0x00), ORANGE,
             label, font_size=7, font_color=WHITE)

add_textbox(slide, 7.0, 4.4, 5.8, 1.5,
            "Flow: DU/CU → UDP:55555 → metrics-server → influxdb → grafana\n\n"
            "All DU/CU containers are dual-homed:\n"
            "  eth0 = oran-intel (priority: 1000)\n"
            "  eth1 = telemetry-net (priority: 100)\n\n"
            "Metrics-only containers (metrics-server, influxdb, grafana)\n"
            "are on telemetry-net ONLY — isolated from data plane.",
            font_size=11, color=GRAY)

# =========================================================================
# SLIDE 3: IP Address Table
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.3, 0.2, 12, 0.6,
            "Complete IP Address Assignment Table",
            font_size=24, color=BLUE, bold=True)

table_text = """
Container        oran-intel (eth0)         telemetry-net (eth1)        Role
─────────        ─────────────────         ──────────────────        ──────────────────
5GC              175.53.1.2                -                         5G Core (AMF/UPF)
cu0              175.53.10.1 + .1.11*      171.53.10.1               CU (gnb_id: 511)
cu1              175.53.10.2 + .1.12*      171.53.10.2               CU (gnb_id: 512)
cu2              175.53.10.3 + .1.15*      171.53.10.3               CU (gnb_id: 513)
du0              175.53.1.10               171.53.1.10               DU → cu0
du1              175.53.1.13               171.53.1.13               DU → cu0
du2              175.53.1.14               171.53.1.14               DU → cu1
du3              175.53.1.16               171.53.1.16               DU → cu2
du4              175.53.1.17               171.53.1.17               DU → cu2
du5              175.53.1.18               171.53.1.18               DU → cu2
ue0              175.53.2.1                -                         UE → du0
ue1              175.53.3.1                -                         UE → du1
ue2              175.53.4.1                -                         UE → du2
ue3              175.53.5.1                -                         UE → du3
ue4              175.53.6.1                -                         UE → du4
ue5              175.53.7.1                -                         UE → du5
redis            175.24.0.1                -                         Cache
metrics-server   -                         171.40.1.4                Metric collector
influxdb         -                         171.40.1.5                Time-series DB
grafana          -                         171.40.1.6                Dashboard

* = secondary IP added via entrypoint.sh (dynamic interface detection)
"""

add_textbox(slide, 0.3, 1.0, 12.5, 6.0, table_text,
            font_size=11, color=WHITE)

# =========================================================================
# SLIDE 4: Interface Assignment & tc Packet Loss
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.3, 0.2, 12, 0.6,
            "Container Interface Assignment & tc Packet Loss",
            font_size=24, color=BLUE, bold=True)

iface_text = """
  DUAL-HOMED CONTAINER (e.g., srsdu0)
  ┌─────────────────────────────────────────────────────────┐
  │                                                         │
  │   eth0 (oran-intel)              eth1 (telemetry-net)     │
  │   IP: 175.53.1.10               IP: 171.53.1.10        │
  │   priority: 1000 (HIGH)         priority: 100 (LOW)    │
  │       │                               │                 │
  │   ┌───┴────────────┐           ┌──────┴──────────┐     │
  │   │ DATA PLANE     │           │ METRICS          │     │
  │   │                │           │                  │     │
  │   │ • F1AP → CU    │           │ • UDP:55555 →    │     │
  │   │   (SCTP:38472) │           │   metrics-server │     │
  │   │ • F1U → CU     │           │   (171.40.1.4)   │     │
  │   │   (UDP)        │           │                  │     │
  │   │ • ZMQ → UE     │           │ JSON metrics     │     │
  │   │   (TCP)        │           │ every 500ms      │     │
  │   └────────────────┘           └──────────────────┘     │
  │                                                         │
  │   tc qdisc replace dev eth0   ← CORRECT TARGET         │
  │   root netem loss X%            (affects data plane)    │
  │                                                         │
  └─────────────────────────────────────────────────────────┘

  WHY priority MATTERS:
  ─────────────────────
  Docker assigns interfaces ALPHABETICALLY by network name:
    "telemetry-net" (m) < "oran-intel" (o) → eth0 would be telemetry-net!

  Fix: priority: 1000 on oran-intel forces it to be eth0
       priority: 100  on telemetry-net  makes it eth1

  This ensures:
    • tc on eth0 = affects srsRAN data plane traffic
    • Metrics flow on eth1 = unaffected by tc experiments
"""

add_textbox(slide, 0.3, 1.0, 12.5, 6.0, iface_text,
            font_size=11, color=WHITE)


# =========================================================================
# SLIDE 5: Why tc Doesn't Trigger HARQ in ZMQ Mode
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.3, 0.2, 12, 0.6,
            "Why tc Packet Loss Does NOT Trigger HARQ in ZMQ Mode",
            font_size=24, color=RED, bold=True)

harq_text = """
  REAL 5G (over-the-air):                ZMQ SIMULATION (this setup):

  UE         DU                          UE              DU
  ┌────┐   ┌────┐                        ┌────┐        ┌────┐
  │ APP│   │    │                        │ APP│        │    │
  │ TCP│   │    │                        │ TCP│        │    │
  │ IP │   │    │                        │ IP │        │    │
  │PDCP│   │PDCP│                        │PDCP│        │PDCP│
  │RLC │   │RLC │                        │RLC │        │RLC │
  │MAC │◄─►│MAC │ ← HARQ here           │MAC │◄──────►│MAC │ ← HARQ here
  │PHY │◄─►│PHY │                        │PHY │        │PHY │
  └────┘   └────┘                        └──┬─┘        └──┬─┘
     ▲ radio  ▲                              │  ZMQ(TCP)   │
     └────────┘                              └──────┬──────┘
   noise → CRC fail                                 │
   → HARQ NACK                               ┌─────┴──────┐
   → retransmit                               │ Linux net  │ ← tc operates HERE
                                              │ (IP/TCP)   │
                                              └────────────┘

  PROBLEM:  tc drops IP packet
                  ↓
            TCP detects loss, retransmits SILENTLY
                  ↓
            ZMQ receives complete data
                  ↓
            PHY gets PERFECT IQ samples
                  ↓
            MAC decodes perfectly → NO HARQ needed

  RESULT:   HARQ = 0 always in ZMQ mode regardless of tc loss
            Bitrate keeps INCREASING (scheduler sees 0% BLER → raises MCS)

  TO TRIGGER HARQ: Need PHY-level channel model (AWGN noise), not tc
"""

add_textbox(slide, 0.3, 1.0, 12.5, 6.2, harq_text,
            font_size=11, color=WHITE)


# =========================================================================
# SLIDE 6: What tc CAN Affect (F1U)
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.3, 0.2, 12, 0.6,
            "What tc CAN Affect: F1U (DU ↔ CU) User Plane",
            font_size=24, color=GREEN, bold=True)

f1u_text = """
  When tc is applied on eth0 (oran-intel) of a DU container:

  ┌────────┐         tc netem loss X%          ┌────────┐
  │  DU    │──────── eth0 (175.x.x.x) ────────│  CU    │
  │        │         F1U = UDP (no retx!)      │        │
  │        │         F1AP = SCTP (has retx)    │        │
  └────────┘                                   └────────┘

  EFFECT ON F1U (UDP):
  ────────────────────
  • F1U uses UDP → NO automatic retransmission
  • Dropped F1U packets = lost user plane data
  • RLC layer detects gaps → RLC retransmissions increase
  • PDCP may also trigger retransmissions
  • Overall throughput/bitrate DECREASES

  EFFECT ON F1AP (SCTP):
  ──────────────────────
  • F1AP uses SCTP → HAS retransmission
  • Control plane messages survive tc loss (with delay)
  • DU-CU connection stays alive

  EFFECT ON ZMQ (DU ↔ UE):
  ─────────────────────────
  • ZMQ uses TCP → HAS retransmission
  • Air interface samples survive tc loss
  • HARQ still won't trigger (PHY sees perfect data)

  CONCLUSION:
  ───────────
  tc on eth0 will cause:
    ✓ Bitrate drop (F1U packet loss)
    ✓ RLC retransmissions increase
    ✗ HARQ retransmissions (still zero - ZMQ is perfect)

  Metrics collection on eth1 remains UNAFFECTED.
"""

add_textbox(slide, 0.3, 1.0, 12.5, 6.2, f1u_text,
            font_size=12, color=WHITE)


# =========================================================================
# SLIDE 7: Compose File Structure
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.3, 0.2, 12, 0.6,
            "Docker Compose File Structure",
            font_size=24, color=BLUE, bold=True)

compose_text = """
  ┌─────────────────────────────────────────────────────────────────────┐
  │ docker-compose-cu++.yaml                                            │
  │   • 5GC (open5gs)     - oran-intel only                            │
  │   • cu0, cu1, cu2     - oran-intel (pri:1000) + telemetry-net (pri:100)│
  │   • entrypoint.sh     - adds secondary 175.x IP (dynamic iface)   │
  └─────────────────────────────────────────────────────────────────────┘

  ┌─────────────────────────────────────────────────────────────────────┐
  │ docker-compose-du.yaml                                              │
  │   • du0-du5           - oran-intel (pri:1000) + telemetry-net (pri:100)│
  │   • redis             - oran-intel only                            │
  └─────────────────────────────────────────────────────────────────────┘

  ┌─────────────────────────────────────────────────────────────────────┐
  │ docker-compose-ue++.yaml                                            │
  │   • ue0-ue5           - oran-intel only                            │
  │   • metrics-server    - telemetry-net only (171.40.1.4)              │
  │   • influxdb          - telemetry-net only (171.40.1.5)              │
  │   • grafana           - telemetry-net only (171.40.1.6)              │
  │   • redis             - oran-intel only                            │
  └─────────────────────────────────────────────────────────────────────┘

  Config files changed:
  ─────────────────────
  • RAN/du_zmq*.conf    → metrics addr: 171.40.1.4 (was 175.40.1.4)
  • RAN/cu_*.yml        → metrics addr: 171.40.1.4 (was 175.40.1.4)
  • RAN/entrypoint_*.sh → dynamic iface detection (was hardcoded eth0)
"""

add_textbox(slide, 0.3, 1.0, 12.5, 6.0, compose_text,
            font_size=12, color=WHITE)


# Save
prs.save(OUTPUT_FILE)
print(f"PPT saved to: {OUTPUT_FILE}")
