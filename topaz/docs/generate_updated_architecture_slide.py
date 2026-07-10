#!/usr/bin/env python3
"""
Generate Updated Architecture Slide for Topology_Agnostic_AD_Presentation.pptx
===============================================================================
Replaces/adds a slide showing:
  - CU-DU-UE topology with ARROWS for connections
  - M-plane (telemetry-net) metric flow
  - Three planes labeled: C-Plane, U-Plane, M-Plane
  - tc packet loss injection point

Usage:
    python3 docs/generate_updated_architecture_slide.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import copy

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
INPUT_FILE = os.path.join(BASE_DIR, "Topology_Agnostic_AD_Presentation.pptx")
OUTPUT_FILE = os.path.join(BASE_DIR, "Topology_Agnostic_AD_Presentation.pptx")

# Color palette
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x16, 0x21, 0x3E)
BLUE = RGBColor(0x00, 0xD2, 0xFF)
GREEN = RGBColor(0x00, 0xFF, 0x88)
RED = RGBColor(0xFF, 0x44, 0x44)
ORANGE = RGBColor(0xFF, 0xAA, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE = RGBColor(0xBB, 0x86, 0xFC)
YELLOW = RGBColor(0xFF, 0xEB, 0x3B)
DARK_BLUE = RGBColor(0x00, 0x22, 0x55)
DARK_GREEN = RGBColor(0x00, 0x44, 0x00)
DARK_ORANGE = RGBColor(0x44, 0x22, 0x00)
DARK_PURPLE = RGBColor(0x33, 0x22, 0x55)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=12,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Consolas"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None,
             text="", font_size=10, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = "Consolas"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     border_color=None, text="", font_size=10,
                     font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = "Consolas"
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top,
              color=WHITE, width=Pt(2)):
    """Add a connector arrow from (start) to (end) in inches."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(start_left), Inches(start_top),
        Inches(end_left), Inches(end_top)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead at end
    connector.end_x = Inches(end_left)
    connector.end_y = Inches(end_top)
    return connector


def add_line(slide, x1, y1, x2, y2, color=WHITE, width_pt=1.5, dashed=False):
    """Add a line shape from (x1,y1) to (x2,y2) in inches."""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    if dashed:
        from pptx.oxml.ns import qn
        ln = connector.line._ln
        prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(prstDash)
    return connector


def build_architecture_slide(prs):
    """Build the full architecture slide with arrows."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_DARK)

    # ---- TITLE BAR ----
    add_rect(slide, 0, 0, 13.333, 0.8, RGBColor(0x0A, 0x0A, 0x1E), BLUE)
    add_textbox(slide, 0.3, 0.1, 9, 0.6,
                "Lab Architecture: CU-DU-UE Topology (Data + Metrics + Control Planes)",
                font_size=20, color=BLUE, bold=True)

    # ---- LEGEND (top-right) ----
    add_rect(slide, 9.8, 0.05, 3.3, 0.7, BG_CARD, GRAY)
    add_textbox(slide, 9.9, 0.05, 3, 0.18,
                "C-Plane (F1AP/N2 SCTP)", font_size=8, color=BLUE)
    add_textbox(slide, 9.9, 0.22, 3, 0.18,
                "U-Plane (F1U UDP / ZMQ TCP)", font_size=8, color=GREEN)
    add_textbox(slide, 9.9, 0.39, 3, 0.18,
                "M-Plane (UDP:55555 telemetry-net)", font_size=8, color=ORANGE)
    add_textbox(slide, 9.9, 0.55, 3, 0.18,
                "tc packet loss injection", font_size=8, color=RED)

    # ==================================================================
    # ROW 1: 5GC (top center)
    # ==================================================================
    gc_left, gc_top = 5.5, 0.95
    add_rounded_rect(slide, gc_left, gc_top, 2.2, 0.6, DARK_GREEN, GREEN,
                     "5GC (AMF/UPF)\n175.53.1.2", font_size=9, font_color=GREEN)

    # ==================================================================
    # ROW 2: CUs
    # ==================================================================
    cu_y = 2.0
    cu_w, cu_h = 2.0, 0.65

    # CU0 (du0, du1)
    cu0_x = 1.0
    add_rounded_rect(slide, cu0_x, cu_y, cu_w, cu_h, DARK_BLUE, BLUE,
                     "CU0 (gnb:511)\n175.53.10.1", font_size=8, font_color=BLUE)

    # CU1 (du2)
    cu1_x = 5.5
    add_rounded_rect(slide, cu1_x, cu_y, cu_w, cu_h, DARK_BLUE, BLUE,
                     "CU1 (gnb:512)\n175.53.10.2", font_size=8, font_color=BLUE)

    # CU2 (du3, du4, du5)
    cu2_x = 9.8
    add_rounded_rect(slide, cu2_x, cu_y, cu_w, cu_h, DARK_BLUE, BLUE,
                     "CU2 (gnb:513)\n175.53.10.3", font_size=8, font_color=BLUE)

    # ---- 5GC <-> CU arrows (N2/SCTP = C-Plane BLUE, dashed) ----
    # 5GC center bottom -> CU top center
    gc_cx = gc_left + 1.1
    gc_bot = gc_top + 0.6

    for cu_x in [cu0_x, cu1_x, cu2_x]:
        cu_cx = cu_x + cu_w / 2
        # N2 control (blue dashed)
        add_line(slide, gc_cx, gc_bot, cu_cx, cu_y, BLUE, 1.5, dashed=True)
        # N3 user plane (green solid) - slightly offset
        add_line(slide, gc_cx + 0.08, gc_bot, cu_cx + 0.08, cu_y, GREEN, 1.5)

    # Label the connections
    add_textbox(slide, 3.0, 1.5, 1.5, 0.25, "N2(SCTP)", font_size=7, color=BLUE)
    add_textbox(slide, 7.2, 1.5, 1.5, 0.25, "N3(GTP-U)", font_size=7, color=GREEN)

    # ==================================================================
    # ROW 3: DUs
    # ==================================================================
    du_y = 3.2
    du_w, du_h = 1.5, 0.65

    du_positions = {
        "du0": (0.3,  "175.53.1.10", "cu0"),
        "du1": (2.0,  "175.53.1.13", "cu0"),
        "du2": (5.2,  "175.53.1.14", "cu1"),
        "du3": (8.2,  "175.53.1.16", "cu2"),
        "du4": (9.9,  "175.53.1.17", "cu2"),
        "du5": (11.5, "175.53.1.18", "cu2"),
    }

    cu_centers = {
        "cu0": cu0_x + cu_w / 2,
        "cu1": cu1_x + cu_w / 2,
        "cu2": cu2_x + cu_w / 2,
    }
    cu_bot = cu_y + cu_h

    for du_name, (dx, ip, cu_parent) in du_positions.items():
        # DU box
        border = RED if du_name == "du0" else BLUE
        add_rounded_rect(slide, dx, du_y, du_w, du_h,
                         RGBColor(0x1A, 0x2A, 0x44), border,
                         f"{du_name}\n{ip}", font_size=8, font_color=WHITE)

        # tc loss marker on du0
        if du_name == "du0":
            add_textbox(slide, dx + 0.05, du_y - 0.22, 1.4, 0.22,
                        "tc netem loss X%", font_size=7, color=RED, bold=True)

        du_cx = dx + du_w / 2
        cu_cx = cu_centers[cu_parent]

        # F1AP (SCTP) = C-Plane blue dashed
        add_line(slide, cu_cx - 0.05, cu_bot, du_cx - 0.05, du_y, BLUE, 1.2, dashed=True)
        # F1U (UDP) = U-Plane green solid
        add_line(slide, cu_cx + 0.05, cu_bot, du_cx + 0.05, du_y, GREEN, 1.2)

    # Labels between CU-DU
    add_textbox(slide, 0.3, 2.75, 1.5, 0.2, "F1AP(SCTP:38472)", font_size=6, color=BLUE)
    add_textbox(slide, 2.0, 2.75, 1.3, 0.2, "F1U(UDP)", font_size=6, color=GREEN)

    # ==================================================================
    # ROW 4: UEs
    # ==================================================================
    ue_y = 4.4
    ue_w, ue_h = 1.5, 0.55

    ue_data = {
        "ue0": (0.3,  "175.53.2.1", "du0"),
        "ue1": (2.0,  "175.53.3.1", "du1"),
        "ue2": (5.2,  "175.53.4.1", "du2"),
        "ue3": (8.2,  "175.53.5.1", "du3"),
        "ue4": (9.9,  "175.53.6.1", "du4"),
        "ue5": (11.5, "175.53.7.1", "du5"),
    }

    du_bot = du_y + du_h

    for ue_name, (ux, ip, du_parent) in ue_data.items():
        add_rounded_rect(slide, ux, ue_y, ue_w, ue_h, DARK_PURPLE, PURPLE,
                         f"{ue_name}\n{ip}", font_size=8, font_color=WHITE)

        ue_cx = ux + ue_w / 2
        du_px = du_positions[du_parent][0] + du_w / 2

        # ZMQ (TCP) = green solid (user data)
        add_line(slide, du_px, du_bot, ue_cx, ue_y, GREEN, 1.2)

    add_textbox(slide, 5.2, 4.1, 1.5, 0.2, "ZMQ(TCP)", font_size=6, color=GREEN)

    # ==================================================================
    # M-PLANE (telemetry-net) - RIGHT SIDE
    # ==================================================================
    # Metrics infrastructure column
    mp_x = 3.8
    mp_y = 5.2

    # M-plane label bar
    add_rect(slide, 0.2, 5.05, 13.0, 0.25, RGBColor(0x33, 0x22, 0x00), ORANGE)
    add_textbox(slide, 0.3, 5.05, 12, 0.25,
                "M-PLANE (telemetry-net 171.0.0.0/8) -- All DUs & CUs send metrics via eth1",
                font_size=9, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    # Metrics server
    ms_x, ms_y = 1.5, 5.5
    add_rounded_rect(slide, ms_x, ms_y, 2.2, 0.55, DARK_ORANGE, ORANGE,
                     "metrics-server\n171.40.1.4:55555", font_size=8, font_color=ORANGE)

    # InfluxDB
    idb_x, idb_y = 5.0, 5.5
    add_rounded_rect(slide, idb_x, idb_y, 1.8, 0.55, DARK_ORANGE, ORANGE,
                     "InfluxDB\n171.40.1.5:8086", font_size=8, font_color=ORANGE)

    # Grafana
    gf_x, gf_y = 8.0, 5.5
    add_rounded_rect(slide, gf_x, gf_y, 1.8, 0.55, DARK_ORANGE, ORANGE,
                     "Grafana\n171.40.1.6:3000", font_size=8, font_color=ORANGE)

    # M-plane flow arrows
    add_line(slide, ms_x + 2.2, ms_y + 0.28, idb_x, idb_y + 0.28, ORANGE, 2.0)
    add_line(slide, idb_x + 1.8, idb_y + 0.28, gf_x, gf_y + 0.28, ORANGE, 2.0)

    add_textbox(slide, 3.5, 5.35, 1.5, 0.18, "UDP:55555", font_size=7, color=ORANGE)
    add_textbox(slide, 7.0, 5.35, 1.0, 0.18, "HTTP", font_size=7, color=ORANGE)

    # Dashed orange lines from each DU/CU down to M-plane bar
    for du_name, (dx, ip, _) in du_positions.items():
        du_cx = dx + du_w / 2
        add_line(slide, du_cx, du_y + du_h, du_cx, 5.05, ORANGE, 0.8, dashed=True)

    for cu_cx_val in [cu0_x + cu_w/2, cu1_x + cu_w/2, cu2_x + cu_w/2]:
        add_line(slide, cu_cx_val, cu_y + cu_h, cu_cx_val, 5.05, ORANGE, 0.8, dashed=True)

    # ==================================================================
    # BOTTOM: Key info
    # ==================================================================
    add_textbox(slide, 0.3, 6.2, 6, 0.5,
                "eth0 = oran-intel (pri:1000) -- Data+Control\n"
                "eth1 = telemetry-net (pri:100) -- Metrics only",
                font_size=9, color=GRAY)
    add_textbox(slide, 6.5, 6.2, 6, 0.5,
                "tc on eth0 -> F1U(UDP) loss -> bitrate drop, RLC retx\n"
                "ZMQ(TCP) retransmits -> HARQ=0 always in ZMQ mode",
                font_size=9, color=GRAY)

    add_textbox(slide, 0.3, 6.7, 12.5, 0.5,
                "Redis: 175.24.0.1 | Metrics flow: DU/CU --[eth1 UDP:55555]--> metrics-server --> InfluxDB --> Grafana",
                font_size=9, color=LIGHT_GRAY)

    return slide


def move_slide_to_index(prs, from_idx, to_idx):
    """Move a slide from one position to another using XML manipulation."""
    from lxml import etree
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    el = slides[from_idx]
    xml_slides.remove(el)
    # Re-read after removal
    slides = list(xml_slides)
    if to_idx >= len(slides):
        xml_slides.append(el)
    else:
        slides[to_idx].addprevious(el)


def delete_slide(prs, idx):
    """Delete a slide by index."""
    from lxml import etree
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[idx].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[idx])


def main():
    prs = Presentation(INPUT_FILE)

    # Find slide 39 (0-indexed = 38) and replace it
    # Slide 39 is "Lab Architecture: CU-DU-UE Topology with Packet Loss Location"
    target_idx = 38
    total_before = len(prs.slides)
    print(f"Total slides before: {total_before}")

    # Build the new slide (appended at end)
    new_slide = build_architecture_slide(prs)
    new_idx = len(prs.slides) - 1
    print(f"New slide added at index {new_idx}")

    # Delete old slide 39
    delete_slide(prs, target_idx)
    print(f"Deleted old slide at index {target_idx}")

    # Now the new slide is at (new_idx - 1) since we removed one before it
    current_new_idx = new_idx - 1
    # Move it to position 38
    move_slide_to_index(prs, current_new_idx, target_idx)
    print(f"Moved new slide to index {target_idx}")

    prs.save(OUTPUT_FILE)
    print(f"Updated architecture slide (slide 39) saved to: {OUTPUT_FILE}")
    print(f"Total slides after: {len(prs.slides)}")


if __name__ == "__main__":
    main()
